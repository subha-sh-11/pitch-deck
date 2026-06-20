"""Parse an uploaded reference deck (.pptx) into a style + structure reference.

Used by the "reference deck" feature: the director uploads an existing deck and the
generator mimics its look (colours, fonts) and structure (slide sequence + titles).
Reads with python-pptx (already a dependency for export). Never raises on bad input.
"""
from __future__ import annotations

import io
import re

from app.core.logging import get_logger

_log = get_logger("pptx_ref")

# Words too common to signal a real title match between an outline slide and a reference slide.
_STOP = {"the", "a", "an", "of", "and", "to", "in", "on", "for", "with", "our",
         "slide", "page"}


def _tokens(text: str) -> set[str]:
    return {w for w in re.findall(r"[a-z0-9]+", (text or "").lower())
            if w not in _STOP and len(w) > 1}


def match_slide(title: str, purpose: str, reference: dict | None) -> dict | None:
    """Best-matching reference slide for an outline slide, by title/purpose word overlap.

    Returns the reference slide ({title, text}) when there's a meaningful overlap, else None
    (so we never feed the writer an unrelated reference slide). Used to ground each generated
    slide in the framing of the corresponding slide of the uploaded reference deck.
    """
    slides = (reference or {}).get("slides") or []
    if not slides:
        return None
    want = _tokens(title) | _tokens(purpose)
    if not want:
        return None
    best, best_score = None, 0
    for sl in slides:
        if not isinstance(sl, dict):
            continue
        score = len(want & (_tokens(sl.get("title")) | _tokens(sl.get("text"))))
        if score > best_score:
            best, best_score = sl, score
    return best if best_score >= 1 else None


def _font_color_hex(run) -> str | None:
    """Best-effort explicit RGB hex for a run's font color (None for theme/auto)."""
    try:
        color = run.font.color
        if color is not None and color.type is not None and color.rgb is not None:
            return f"#{str(color.rgb)}"
    except Exception:
        return None
    return None


def extract_reference(data: bytes) -> dict:
    """Return {slideCount, slides:[{title,text}], fonts:[...], colors:[#hex]}."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides: list[dict] = []
    fonts: dict[str, int] = {}
    colors: dict[str, int] = {}

    for slide in prs.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    texts.append(line)
                for run in para.runs:
                    name = getattr(run.font, "name", None)
                    if name:
                        fonts[name] = fonts.get(name, 0) + 1
                    hexc = _font_color_hex(run)
                    if hexc:
                        colors[hexc] = colors.get(hexc, 0) + 1

        title = ""
        try:
            if slide.shapes.title is not None and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
        except Exception:
            title = ""
        if not title and texts:
            title = texts[0]

        slides.append({"title": title[:160], "text": " ".join(texts)[:600]})

    top_fonts = [f for f, _ in sorted(fonts.items(), key=lambda kv: -kv[1])[:3]]
    top_colors = [c for c, _ in sorted(colors.items(), key=lambda kv: -kv[1])[:6]]
    _log.info("pptx_ref: %d slides, %d fonts, %d colors", len(slides), len(top_fonts), len(top_colors))
    return {
        "slideCount": len(slides),
        "slides": slides[:40],
        "fonts": top_fonts,
        "colors": top_colors,
    }
