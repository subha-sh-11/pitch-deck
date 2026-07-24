"""Turn an uploaded reference DECK (.pptx / PDF) into IMAGES for the vision model.

The parsed reference deck (`pptx_ref.extract_reference`) only captures text, fonts and
explicit run colours — the vision-based reference-analysis stage never SEES the deck's
composition. This module closes that gap:

- PDF decks (Canva exports are usually PDF): render the first few pages to JPEG.
- .pptx decks: no pure-python slide renderer exists, so extract the embedded pictures
  (the mood/still images the deck is built from) instead.

The router persists the results as ``upload_ref`` assets tagged
``generation_meta.source='visual_direction'`` with names prefixed ``refdeck-`` — the
existing generation pipeline (``_load_references_full`` → reference_analysis agent) then
picks them up with no changes to the generation flow.

Never raises on bad input: every helper returns [] on failure (callers treat deck images
as best-effort enrichment, not a hard requirement).
"""
from __future__ import annotations

import io
import re

from app.core.logging import get_logger

_log = get_logger("deck_ref_images")

MAX_PDF_PAGES = 6          # first pages carry the deck's visual identity
MAX_PPTX_IMAGES = 8
MIN_PPTX_IMAGE_PX = 200    # skip icons / logos / decorative slivers
TARGET_WIDTH = 1024        # ~vision-model sweet spot; keeps payloads small
JPEG_QUALITY = 80


def safe_stem(filename: str, max_len: int = 40) -> str:
    """Filename → a safe lowercase stem for asset names ('My Deck (v2).pdf' → 'my-deck-v2')."""
    stem = re.sub(r"\.[a-z0-9]+$", "", (filename or "deck").lower())
    stem = re.sub(r"[^a-z0-9]+", "-", stem).strip("-")
    return (stem or "deck")[:max_len]


def _encode_jpeg(pil_image) -> bytes:
    from PIL import Image  # noqa: F401 — ensure Pillow present before .convert

    img = pil_image.convert("RGB")
    if img.width > TARGET_WIDTH:
        img = img.resize(
            (TARGET_WIDTH, max(1, round(img.height * TARGET_WIDTH / img.width))))
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=JPEG_QUALITY)
    return buf.getvalue()


def render_pdf_pages(data: bytes, max_pages: int = MAX_PDF_PAGES) -> list[bytes]:
    """PDF bytes → JPEG bytes for the first ``max_pages`` pages (~TARGET_WIDTH px wide).

    Uses pypdfium2 (already a dependency — the scanned-script OCR fallback renders with it).
    """
    try:
        import pypdfium2 as pdfium  # lazy — pure pip wheel

        pdf = pdfium.PdfDocument(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001 — corrupt / encrypted / non-PDF
        _log.warning("deck_ref_images: could not open PDF: %s", exc)
        return []
    out: list[bytes] = []
    try:
        for i in range(min(len(pdf), max_pages)):
            try:
                page = pdf[i]
                width = float(page.get_width() or 612.0)
                scale = max(0.1, min(4.0, TARGET_WIDTH / width))
                bitmap = page.render(scale=scale)
                out.append(_encode_jpeg(bitmap.to_pil()))
                page.close()
            except Exception as exc:  # noqa: BLE001 — skip the bad page, keep the rest
                _log.warning("deck_ref_images: PDF page %d render failed: %s", i + 1, exc)
    finally:
        try:
            pdf.close()
        except Exception:  # noqa: BLE001
            pass
    _log.info("deck_ref_images: rendered %d PDF page(s)", len(out))
    return out


def extract_pptx_images(data: bytes, max_images: int = MAX_PPTX_IMAGES) -> list[bytes]:
    """.pptx bytes → JPEG bytes of the pictures embedded in its slides, in slide order.

    Skips tiny images (< MIN_PPTX_IMAGE_PX on either side — icons/logos) and duplicates
    (decks often repeat one background still on every slide).
    """
    try:
        from PIL import Image
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        _log.warning("deck_ref_images: could not open pptx: %s", exc)
        return []
    out: list[bytes] = []
    seen_blobs: set[int] = set()
    for slide in prs.slides:
        if len(out) >= max_images:
            break
        for shape in slide.shapes:
            if len(out) >= max_images:
                break
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                blob = shape.image.blob
            except Exception:  # noqa: BLE001 — linked (non-embedded) or corrupt picture
                continue
            key = hash(blob)
            if key in seen_blobs:
                continue
            seen_blobs.add(key)
            try:
                img = Image.open(io.BytesIO(blob))
                if min(img.width, img.height) < MIN_PPTX_IMAGE_PX:
                    continue
                out.append(_encode_jpeg(img))
            except Exception as exc:  # noqa: BLE001 — WMF/EMF and friends Pillow can't read
                _log.debug("deck_ref_images: skipped unreadable pptx image: %s", exc)
    _log.info("deck_ref_images: extracted %d pptx image(s)", len(out))
    return out


def extract_reference_pdf(data: bytes) -> dict:
    """PDF bytes → the same structure `pptx_ref.extract_reference` yields for .pptx:
    {slideCount, slides:[{title,text}], fonts:[...], colors:[#hex]} — so a PDF reference
    deck also drives outline mirroring / slide matching, not just the page images.
    """
    import pdfplumber  # lazy — already a dependency (script ingestion)

    slides: list[dict] = []
    fonts: dict[str, int] = {}
    colors: dict[str, int] = {}
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages[:40]:
            try:
                text = page.extract_text() or ""
            except Exception:  # noqa: BLE001
                text = ""
            lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
            title = lines[0] if lines else ""
            slides.append({"title": title[:160], "text": " ".join(lines)[:600]})
            try:
                for ch in page.chars[:2000]:
                    name = str(ch.get("fontname") or "")
                    # Strip the subset prefix ('ABCDEF+Helvetica-Bold' → 'Helvetica-Bold').
                    name = name.split("+", 1)[-1]
                    if name:
                        fonts[name] = fonts.get(name, 0) + 1
                    col = ch.get("non_stroking_color")
                    if isinstance(col, (tuple, list)) and len(col) == 3:
                        try:
                            hexc = "#%02X%02X%02X" % tuple(
                                max(0, min(255, round(float(c) * 255))) for c in col)
                            colors[hexc] = colors.get(hexc, 0) + 1
                        except Exception:  # noqa: BLE001
                            pass
            except Exception:  # noqa: BLE001 — char metadata is best-effort
                pass

    top_fonts = [f for f, _ in sorted(fonts.items(), key=lambda kv: -kv[1])[:3]]
    top_colors = [c for c, _ in sorted(colors.items(), key=lambda kv: -kv[1])[:6]]
    _log.info("deck_ref_images: pdf ref — %d pages, %d fonts, %d colors",
              len(slides), len(top_fonts), len(top_colors))
    return {
        "slideCount": len(slides),
        "slides": slides,
        "fonts": top_fonts,
        "colors": top_colors,
    }
